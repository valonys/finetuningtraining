"""PPTX parser — extracts text from slides, notes, and tables."""
from __future__ import annotations

from ..ingest import IngestedRecord


def parse_pptx(path: str, **_kw) -> IngestedRecord:
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError("python-pptx not installed — `pip install python-pptx`") from e

    prs = Presentation(path)
    parts: list[str] = []
    tables: list[list[list[str]]] = []

    for i, slide in enumerate(prs.slides):
        parts.append(f"## Slide {i + 1}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for p in shape.text_frame.paragraphs:
                    txt = p.text.strip()
                    if txt:
                        parts.append(txt)
            if shape.has_table:
                rows: list[list[str]] = []
                for row in shape.table.rows:
                    rows.append([c.text.strip() for c in row.cells])
                tables.append(rows)
                parts.append("\n".join([" | ".join(r) for r in rows]))
        # Speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"_notes:_ {notes}")

    return IngestedRecord(
        source=path,
        source_type="pptx",
        text="\n\n".join(parts),
        tables=tables,
        metadata={"num_slides": len(prs.slides)},
    )
